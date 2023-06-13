from unittest import TestCase
import tempfile
import unittest
#import mock
import urllib.request, urllib.parse, urllib.error
import json
import os
import sys
sys.path.append('..') 
import pagesetparser.core as core
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pprint import pprint
import json
import io
import os
import math
import string
import unicodedata
from PIL import Image
from sys import argv
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE




class ovfTestCommands(TestCase):



    def test_process_commandstring(self):
        button={}
        commandstring='ovf(place("hello"))'
        core.process_commandstring(commandstring,button)
        self.assertEqual(button["vocalization"],"hello")

    def test_process_commandstring_single_quotes(self):
        button={}
        commandstring="ovf(place('hello'))"
        core.process_commandstring(commandstring,button)
        self.assertEqual(button["vocalization"],"hello")

    def test_process_commandstring_single_quotes_with_internal_coma(self):
        button={}
        commandstring="ovf(place('hello there, general kenobi'))"
        core.process_commandstring(commandstring,button)
        self.assertEqual(button["vocalization"],"hello there, general kenobi")


    def test_process_commandstring_with_comma(self):
        button={}
        commandstring='ovf(place(", hello"))'
        core.process_commandstring(commandstring,button)
        self.assertEqual(button["vocalization"],", hello")

    def test_process_commandstring_with_open(self):
        button={}
        commandstring='ovf(open("jazz"))'
        core.process_commandstring(commandstring,button)
        self.assertEqual(button["load_board"]["path"], "boards/jazz.obf")

    def test_process_commandstring_with_two_commands(self):
        button={}
        commandstring='ovf(open("jazz"),place("goodbye"))'
        core.process_commandstring(commandstring,button)
        self.assertEqual(button["load_board"]["path"], "boards/jazz.obf")
        self.assertEqual(button["vocalization"],"goodbye")

    def test_process_commandstring_with_two_commands_without_quotes(self):
        button={}
        commandstring='ovf(open(jazz),place(goodbye))'
        core.process_commandstring(commandstring,button)
        self.assertEqual(button["load_board"]["path"], "boards/jazz.obf")
        self.assertEqual(button["vocalization"],"goodbye")

    def test_what_if_they_forget_the_ovf(self):
        button={}
        commandstring='open(jazz),place(goodbye)'
        core.process_commandstring(commandstring,button)
        self.assertEqual(button["load_board"]["path"], "boards/jazz.obf")
        self.assertEqual(button["vocalization"],"goodbye")

#   This commenting is mentioned in the main file - it's for if we end up with commands that don't need parenthesis 
#    def test_process_commandstring_with_single_wordss(self):
#        button={}
#        commandstring='ovf(alpha,bravo)'
#        core.process_commandstring(commandstring,button)
#
#    def test_process_commandstring_with_single_wordss_without_ovf(self):
#        button={}
#        commandstring='open,place'
#        core.process_commandstring(commandstring,button)

    def test_process_commandstring_with_double_places(self):
        button={}
        commandstring='ovf(place("jazz"),place("goodbye"))'
        core.process_commandstring(commandstring,button)
        self.assertEqual(button["vocalization"],"jazzgoodbye")

    def test_issue_141(self):
        button={}
        commandstring='ovf(place(abc),open(ABC))' #the lack of the speech marks is likely the thing
        core.process_commandstring(commandstring,button)
        self.assertEqual(button["vocalization"],"abc")

if __name__=="__main__":
    unittest.main()
