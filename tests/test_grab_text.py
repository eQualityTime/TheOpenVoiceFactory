from unittest import TestCase
import tempfile
#TODO rename this file
import unittest
#import mock
import urllib.request, urllib.parse, urllib.error
import json
import os
import sys
sys.path.append('..') 
import pagesetparser.pagesetparser as pagesetparser
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pprint import pprint
import json
import io
import os
import math
import string
import unicodedata
from PIL import Image
from sys import argv
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE




class ovfTest(TestCase):

    CK20=None #This is the pageset

    def get_singleton_CK20(self):
        if not self.CK20:
           # prs = Presentation("tests/CK20V2cutdown.pptx")
            pagesetparser.gridSize=5
            self.CK20 = pagesetparser.Pageset("tests/testinputs/CK20V2cutdown.pptx",pagesetparser.gridSize, False) #TODO - this makes the file hard to run from outside the base directory 
            
        return self.CK20

    def test_first(self):
         self.assertEqual(3,3)


    def test_read_CK20_and_count_slides(self):
        grids = self.get_singleton_CK20().grids
        self.assertEqual(len(grids),10)



    def test_make_title(self):
        tag = pagesetparser.make_title("A sentance")
        self.assertEqual(tag,"asentance")

    def test_make_title_2(self):
        tag = pagesetparser.make_title("ThInGs In MiXeD cAsE")
        self.assertEqual(tag,"thingsinmixedcase")

    def test_make_title_3(self):
        tag = pagesetparser.make_title("Sysbols#")
        self.assertEqual(tag,"sysbols")




if __name__=="__main__":
    unittest.main()
