from unittest import TestCase
import filecmp
import unittest
from tests.regressions.generate_regressions import generate_obz
#import mock
import urllib.request, urllib.parse, urllib.error
import json
import os
import sys
sys.path.append('..') #todo - test removal 
from pagesetparser.pageset import Pageset
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pprint import pprint
import json
import io
import os
import math
import string
import glob
import unicodedata
from PIL import Image
from sys import argv
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE

import logging
logging.basicConfig(level=logging.ERROR)
logger = logging.getLogger(__name__)



class ovfTest(TestCase):



    def test_obz_pages(self):

        self.assertTrue(generate_obz_and_compare_to_records("tests/regressions/regression_tests_size_4/CK12+V2.pptx",4))
        self.assertTrue(generate_obz_and_compare_to_records("tests/regressions/regression_tests_size_5/CK20V2.pptx",5))
        self.assertTrue(generate_obz_and_compare_to_records("tests/regressions/regression_tests_size_5/CK20V2_bg_regression.pptx",5))
        self.assertTrue(generate_obz_and_compare_to_records("tests/regressions/regression_tests_size_5/CK20V2_gr_regression.pptx",5))
        self.assertTrue(generate_obz_and_compare_to_records("tests/regressions/regression_tests_size_5/CK20V2_sk_regression.pptx",5))
        self.assertTrue(generate_obz_and_compare_to_records("tests/regressions/regression_tests_size_5/CK20_slo_regession.pptx",5))
        self.assertTrue(generate_obz_and_compare_to_records("tests/regressions/regression_tests_size_5/CK20V2-place.pptx",5))
        self.assertTrue(generate_obz_and_compare_to_records("tests/regressions/regression_tests_size_5/CK20V2_es_regression.pptx",5))
        self.assertTrue(generate_obz_and_compare_to_records("tests/regressions/regression_tests_size_5/CK20V2_ara_regession.pptx",5))
    

def generate_obz_and_compare_to_records(filename, page_size):
        generate_obz(filename,page_size)
        logger.info("We are in directory {}".format(filename))
        is_same=True # We use this boolean so we get a list of the files that fail the match rather than only the first.
        if not compare(filename+".obz.attempt/manifest.json"):
            logger.error("XXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXX")
            logger.error("Regression Test error on file manifest")
            logger.error("try vimdiff "+ filename+".obz.attempt/manifest.json "+ filename+".obz.target/manifest.json")
            logger.error("XXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXX")
            is_same=False
        logger.info("Now we start the comparisons") 
        for file in glob.glob(filename+".obz.attempt/boards/*.obf"):
            if not compare(file):
                logger.error("XXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXX")
                logger.error("Regression Test error on file {}".format(file))
                logger.error("try vimdiff  {} {}".format(file, file.replace('attempt','target')))
                logger.error("XXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXX")
                is_same=False
        return is_same 

def compare(file):
    target=file.replace('attempt','target')
    return filecmp.cmp(file,target)



if __name__=="__main__":
    unittest.main()
