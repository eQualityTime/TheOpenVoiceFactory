from unittest import TestCase
import urllib.request, urllib.parse, urllib.error
import json
import glob
import os
import sys
sys.path.append('../..')
import pagesetparser
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pprint import pprint
import json
import io
import os
import math
import string
import unicodedata
from PIL import Image
from sys import argv
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
import hashlib
from pagesetparser.pageset import Pageset



# switch the 'attempt' to 'target' to actually genearte them 
def generate_obz(filename,gridSize): 
        gridSize = 5
        obf_dest=filename+".obz.attempt/"
        if not os.path.exists(obf_dest):
            os.makedirs(obf_dest)
        if not os.path.exists(obf_dest+"boards/"):
            os.makedirs(obf_dest+"boards/")
        pageset = Pageset(filename, gridSize)
        pageset.write_to_obz(obf_dest)

def generate(filename, gridSize):
        Pageset(filename,"",gridSize,False).write_json(filename+".json")


if __name__=="__main__":

 for file in glob.glob("regression_tests_size_5/*.pptx"):
        print(file)
        generate_obz(file,5)
 for file in glob.glob("regression_tests_size_4/*.pptx"):
        print(file)
        generate_obz(file,4)




