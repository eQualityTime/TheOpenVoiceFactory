from unittest import TestCase
import unittest
import mock
import grab_text
import urllib
import json
import os
import sys
sys.path.append('/home/joereddington/')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pprint import pprint
import json
import io
import os
import math
import string
import unicodedata
from PIL import Image
from sys import argv
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE




class ovfTest(TestCase):

    CK20=None

    def get_singleton_CK20(self):
        if not self.CK20:
            prs = Presentation("testinputs/CK20V2cutdown.pptx")
            grab_text.gridSize=5
            self.CK20 = grab_text.extract_grid(prs)
        return self.CK20

    def test_first(self):
         self.assertEqual(3,3)



    def test_read_CK20_and_count_slides(self):
        grids = self.get_singleton_CK20()
        self.assertEqual(len(grids),10)


    def test_read_CK20_and_check_titles(self):
        grids = self.get_singleton_CK20()
        self.assertEqual(grids[0].tag,"Top page")


    def test_read_CK20_and_check_color(self):
        grids = self.get_singleton_CK20()
        self.assertEqual(grids[0].colors[2][2],(255,125,236))

    def test_read_CK20_and_check_link(self):
        grids = self.get_singleton_CK20()
        print "XXXXXX"
        print grids[0].labels[4][1]
        print grids[0].links[4][1]
        print "XXXXXX"
        self.assertNotEqual(grids[0].links[4][1],"")




    def test_read_CK20_and_check_link_neg(self):
        grids = self.get_singleton_CK20()
        self.assertEqual(grids[0].links[1][1],"")

    def test_make_title(self):
        tag = grab_text.make_title("A sentance")
        self.assertEqual(tag,"asentance")

    def test_make_title_2(self):
        tag = grab_text.make_title("ThInGs In MiXeD cAsE")
        self.assertEqual(tag,"thingsinmixedcase")

    def test_make_title_3(self):
        tag = grab_text.make_title("Sysbols#")
        self.assertEqual(tag,"sysbols")

    def test_warning_for_missinglink(self):
        grab_text.addwarning=mock.Mock(return_value=None)
        prs = Presentation("testinputs/CK20V2cutdown.pptx") #Can't use the sinlgeton because it might have been called before the mocking
        grab_text.gridSize=5
        grab_text.extract_grid(prs)
        grab_text.addwarning.assert_called()


if __name__=="__main__":
    unittest.main()
