import math
from core import make_title
import core
import string
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL

warningMissingLinks = True
bordercolor = False
IMAGE_WARNING = False

class Grid:

    """Class representing on n by n grid, complete with utterances, links
    colours, and so on. Currently outputs as javascript, should also
    write to json on it's own mertits"""

    def update_links(self, grids):
        for col in range(self.grid_size):
            for row in range(self.grid_size):
                current = self.links[row][col]
                if "slide" in current:
                    # first extract the number
                    number_string = ''.join(
                        c for c in current
                        if c in string.digits)
                    # Then work out the relevant tag
                    self.links[row][col] = grids[
                        int(number_string) - 1].tag
                    # Remember that slides are numbered from
                    # 1 but grids are numbered from 0

    def __init__(self, pres, slide, gridSize, owningPageset):
        self.slide=slide
        self.pageset = owningPageset
        self.grid_size = gridSize
        self.labels = [
            ["" for x in range(self.grid_size)]
            for x in range(self.grid_size)]
        self.links = [
            ["" for x in range(self.grid_size)]
            for x in range(self.grid_size)]
        self.colors = [
            ["" for x in range(self.grid_size)]
            for x in range(self.grid_size)]
        self.icons = [
            ["" for x in range(self.grid_size)]
            for x in range(self.grid_size)]
        self.tag = "unknown"
        self.slide_width = pres.slide_width
        self.slide_height = pres.slide_height
        self.feedback = []
        for shape in slide.shapes:
            self.process_shape(shape)

    def get_col_row(self, top, left):
        # It doesn't make sense to use width and height, since often the
        # midpoint lies outside the cell, particularly in the horizontal directino
        # (probably because text boxes have a default minimum width)
        # It might be worth having a slight fudge right-and-down to make sure
        # we don't miss items who are just outside the top left of the
        # cell.
        # This method is terrible - it worked fine with python 2.7, and when I upgraded to python3, the division caused problems. I fixed it by putting interger division *back*, which is clearly wrong but... 
     #   print("top: {}, left: {}, height: {}, width: {}".format(top,left,self.slide_height,self.slide_width))
        numberofrows=self.grid_size #forreadability  
        numberofcolumns=self.grid_size #forreadability  
        col = math.floor(
            (numberofcolumns * left // self.slide_width) + 0.5)
        row = math.floor(
            (numberofrows * top // self.slide_height) + 0.5)



        return (int(col), int(row))

    def process_shape(self, shape):
        try:
            if shape.is_placeholder:
                if shape.placeholder_format.idx == 0:
                    self.tag = make_title(shape.text)
            (co, ro) = self.get_col_row(
                    shape.top+shape.height/2, shape.left+shape.width/2)
            if ((co >= self.grid_size) or (ro >= self.grid_size)):
                print("Warning, shape outside of page area on page:{}".format(self.tag))
                return
            # Now - let's find out if there is a link...
            click_action = None
            if shape.shape_type != MSO_SHAPE_TYPE.GROUP:
                click_action = shape.click_action
            if click_action.hyperlink.address is not None:
                self.links[co][
                        ro] = click_action.hyperlink.address
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                if bordercolor:
                    self.colors[co][
                            ro] = shape.line.color.rgb
                else:
                    if shape.fill.type==MSO_FILL.SOLID:
                        if (str(shape.fill.fore_color.type)
                                == "SCHEME (2)"):
                            #self.colors[co][ro] = ( 200, 0, 0)
                            pass
                        else:
                            self.colors[co][
                                    ro] = shape.fill.fore_color.rgb

            if shape.has_text_frame:
                self.process_text_frame(shape, co, ro)

            if(warningMissingLinks):
                if (click_action.hyperlink.address is None):
                    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                        if shape.auto_shape_type == MSO_SHAPE.FOLDED_CORNER:
                            if len(self.links[
                                   co][ro]) < 1:
                                self.pageset.addfeedback("Unknown link at slide: " + self.tag + " link here: [{}] [{}] {} ".format(co, ro, self.labels[co][ro]) + self.links[co][ro])

        except (AttributeError, KeyError, NotImplementedError):
            self.pageset.addfeedback(core.returnException())
            return

    def process_text_frame(self, shape, co, ro):
        text = ""
        if "Yes" in self.labels[co][ro]:
            return
        for paragraph in shape.text_frame.paragraphs:
            text += "".join([run.text.replace("’", "'")
                             for run in paragraph.runs])
        if text != "":
            # add the if shape_type is text box
            self.labels[co][ro] = text.strip()


    def create_image_grid(self):
        images = {}
        for shape in self.slide.shapes:
            if not hasattr(shape, "shape_type"):
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                (co, ro) = self.get_col_row(
                        shape.top+shape.height/2, shape.left+shape.width/2)
                if (co, ro) not in images:
                    images[co, ro] = []
                images[co, ro].append(shape)
            if IMAGE_WARNING:
                for col in range(self.grid_size):
                    for row in range(self.grid_size):
                        if (col, row) not in images:
                            if self.labels[col][
                                    row] is not "":#wait... is this the error. 
                                if self.tag not in self.labels[
                                        col][row]:
                                    print("WARNING: image missing at column {}, row  {} (label: {}) on slide:{}".format(col, row, labels[col][row], self.tag))
        return images

