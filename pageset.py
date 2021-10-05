from grid import Grid
from core import make_title
from pptx import Presentation
from PIL import Image
import io
import os

class Pageset:

    # Todo
    # Should throw error if fileanem doens't produce a proper pageset.

    def __init__(self, filename, dest, gridSize, saveimages=True):
        self.prs = Presentation(filename)
        self.grids = []
        self.gridSize=gridSize
        self.feedback = []
        self.extract_grid()
        self.extract_and_label_images(dest, False) #This has to run to get the names right. TODO - make into a smaller, nicer, naming file 

    def addfeedback(self, feedelement):
        self.feedback.append(feedelement)
        # This is a stub - this is to be used to manage the passing of
        # messages to the user.
        print("Feedback added :{}".format(feedelement))

    def getfeedback(self):
        return self.feedback

    def extract_grid(self):
        for index,slide in enumerate(self.prs.slides):
            index="Slide number "+str(index+1)
            self.grids.append(Grid(self.prs, slide, self.gridSize, self,index))
        for grid in self.grids:
            grid.update_links(self.grids)

    def extract_and_label_images(self, dest, SAVE=True):
        image_slight_number = 0
        for slide in self.prs.slides: #TODO - you'll note that we don't use the slide in the loop...
            export_images(
                self.grids[image_slight_number],
                image_slight_number,
                dest,
                SAVE)
            image_slight_number += 1


def export_images(grid, slide_number, dest_folder, SAVE=True):
    """     Second pass through shapes list finds images and saves them.
    We have to do this separately so it's guaranteed we already know what to
    name the images!"""
    labels = grid.labels
    images = grid.create_image_grid()
    # Compose each icon out of all the images in the grid cell.
    for (x, y) in images:
        # Go through all the images, compute bounding
        # box.
        l = min([shape.left for shape in images[x, y]])
        t = min([shape.top for shape in images[x, y]])
        r = max([shape.left +
                 shape.width for shape in images[x, y]])
        b = max([shape.top +
                 shape.height for shape in images[x, y]])
        # Scale gives us the mapping from image pixels to powerpoint
        # distance units. This depends on the resolution
        # of the images.
        scale = min([shape.width/shape.image.size[0]
                     for shape in images[x, y]])

        # Size of combined image, in actual pixels (not PPTX units)
        # If scales differ between objects, we resize
        # them next
        w = int((r-l)/scale)
        h = int((b-t)/scale)

        try:

            composite = Image.new('RGBA', (w, h))

            # Add all the images together.
            for shape in images[x, y]:
                partBox=ready_for_composite(shape,scale,w,h,l,t)
                composite=Image.alpha_composite(composite,partBox)
            # Crop final image.
            bbox = composite.getbbox()
            composite = composite.crop(bbox)

            # Save!
            grid.icons[x][y] = "icons/" + create_icon_name(x, y, labels, grid.links, slide_number)
            name =                  create_icon_name(x, y, labels, grid.links, slide_number)
            if SAVE:
                # + str(slide_number)
                folder = dest_folder+"/icons/"
                if not os.path.exists(folder):
                    os.makedirs(folder)
                composite.save(folder + "" + name)
        except IOError as e:
            print("Error reading image for {} {} (Code 121)".format(x, y))
            if ("cannot find loader for this WMF file" in str(e)):
                print("Error: it appears that the image in column {} row {} of slide {}, is for Windows only, please change the format of that image".format(x, y, slide_number))
            if ("cannot identify image file" in str(e)):
                print("Error: it appears that the image in column {} row {} of slide {}, is for Windows only, please change the format of that image".format(x, y, slide_number))
        except ValueError:
            print("Error reading image for {} {} (Code 127)".format(x, y))
        except IndexError:
            print("Error reading image for {} {}- it is outside the grid".format(x, y))


def ready_for_composite(shape,scale,w,h,l,t):

    part = Image.open(
        io.BytesIO(
            shape.image.blob))
    part.load()
    part = part.crop(crop_to_shape(shape,part))
    partScale = (shape.width / part.size[0])
    # part.size because it might have been cropped
    part = resizeImage(part, partScale / scale)
    partBox = Image.new('RGBA', (w, h))
    partBox.paste( part, (int((shape.left - l)/scale),int( (shape.top - t)/scale)))
    return partBox

def crop_to_shape(shape,part):
    width = part.size[0]
    height = part.size[1]
    left = shape.crop_left*width
    right = (1-shape.crop_right)*width
    top = shape.crop_top*height
    bottom = (1-shape.crop_bottom)*height
    box = (int(left),
           int(top),
           int(right),
           int(bottom))
    return box



def resizeImage(image, scaleFactor):

    oldSize = image.size
    newSize = (int(scaleFactor*oldSize[0]),
               int(scaleFactor*oldSize[1]))
    return image.resize(newSize, Image.ANTIALIAS)

def create_icon_name(x, y, labels, links, slide_number):
    name = "S"+str(slide_number)+"X"+str(x)+"Y"+str(y)+".png" 
    try:
        name = "S"+str(slide_number)+"X"+str(x)+"Y"+str(y)+make_title(
                labels[x][y])+".png" 
    except IndexError:
        print("Create_icon_name was given an x y that was outside the possible ranges")
    return name
