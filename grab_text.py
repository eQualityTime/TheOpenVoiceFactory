#!/usr/bin/python
# -*- coding: utf-8 -*-
"Extracting Utterances from CommuniKate pagesets designed in PowerPoint"
# Make the images export more effectively
import sys
sys.path.append('/home/joereddington/')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pprint import pprint
import json
import zipfile
import io
import os
import math
import string
import unicodedata
from PIL import Image
from sys import argv
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL


import sys
import linecache
print_exceptions = True
IMAGE_WARNING = False
bordercolor = False

warningMissingLinks = True


def returnException():
    # http://stackoverflow.com/a/20264059
    if print_exceptions is True:
        exc_type, exc_obj, tb = sys.exc_info()
        f = tb.tb_frame
        lineno = tb.tb_lineno
        filename = f.f_code.co_filename
        linecache.checkcache(filename)
        line = linecache.getline(filename, lineno, f.f_globals)
        return 'EXCEPTION IN ({}, LINE {} "{}"): {}'.format(
                filename, lineno, line.strip(), exc_obj)


def resizeImage(image, scaleFactor):

    oldSize = image.size
    newSize = (int(scaleFactor*oldSize[0]),
               int(scaleFactor*oldSize[1]))
    return image.resize(newSize, Image.ANTIALIAS)


def remove_punctuation(s):
    """
    >>> strip_punctuation(u'something')
    u'something'

    >>> strip_punctuation(u'something.,:else really')
    u'somethingelse really'
    """
    text = s
    if not isinstance(s, type("hope")):
        text = str(s, "utf-8")
    punctutation_cats = set(['Pc', 'Pd', 'Ps', 'Pe', 'Pi', 'Pf', 'Po'])
    return ''.join(x for x in text
                   if unicodedata.category(x) not in punctutation_cats)


class Pageset:

    # Todo
    # Should throw error if fileanem doens't produce a proper pageset.

    def __init__(self, filename, dest, saveimages=True):
        self.prs = Presentation(filename)
        self.grids = []
        self.feedback = []
        self.extract_grid()
        self.extract_and_label_images(dest, saveimages)

    def addfeedback(self, feedelement):
        self.feedback.append(feedelement)
        # This is a stub - this is to be used to manage the passing of
        # messages to the user.
        print("Feedback added :{}".format(feedelement))

    def getfeedback(self):
        return self.feedback

    def extract_grid(self):
        for slide in self.prs.slides:
            self.grids.append(Grid(self.prs, slide, gridSize, self))
        for grid in self.grids:
            grid.update_links(self.grids)

    def extract_and_label_images(self, filename, SAVE):
        image_slight_number = 0
        for slide in self.prs.slides:
            export_images(
                self.grids[image_slight_number],
                image_slight_number,
                slide,
                filename,
                SAVE)
            image_slight_number += 1


class Grid:

    """Class representing on n by n grid, complete with utterances, links
    colours, and so on. Currently outputs as javascript, should also
    write to json on it's own mertits"""

    def update_links(self, grids):
        for col in range(self.grid_size):
            for row in range(self.grid_size):
                current = self.links[row][col]
                if "slide" in current:
                    # first extract the number
                    number_string = ''.join(
                        c for c in current
                        if c in string.digits)
                    # Then work out the relevant tag
                    self.links[row][col] = grids[
                        int(number_string) - 1].tag
                    # Remember that slides are numbered from
                    # 1 but grids are numbered from 0

    def __init__(self, pres, slide, gridSize, owningPageset):
        self.grid_size = gridSize
        self.pageset = owningPageset
        self.labels = [
            ["" for x in range(self.grid_size)]
            for x in range(self.grid_size)]
        self.utterances = [
            ["" for x in range(self.grid_size)]
            for x in range(self.grid_size)]
        self.links = [
            ["" for x in range(self.grid_size)]
            for x in range(self.grid_size)]
        self.colors = [
            ["" for x in range(self.grid_size)]
            for x in range(self.grid_size)]
        self.icons = [
            ["" for x in range(self.grid_size)]
            for x in range(self.grid_size)]
        self.tag = "unknown"
        self.slide_width = pres.slide_width
        self.slide_height = pres.slide_height
        self.feedback = []
        for shape in slide.shapes:
            self.process_shape(shape)

    def get_col_row(self, top, left):
        # It doesn't make sense to use width and height, since often the
        # midpoint lies outside the cell, particularly in the horizontal directino
        # (probably because text boxes have a default minimum width)
        # It might be worth having a slight fudge right-and-down to make sure
        # we don't miss items who are just outside the top left of the
        # cell.
        col = math.floor(
            (left * self.grid_size / self.slide_width) + 0.5)
        row = math.floor(
            (top * self.grid_size / self.slide_height) + 0.5)
        return (int(col), int(row))

    def process_shape(self, shape):
        try:
            if shape.is_placeholder:
                if shape.placeholder_format.idx == 0:
                    self.tag = make_title(shape.text)
            (co, ro) = self.get_col_row(
                    shape.top+shape.height/2, shape.left+shape.width/2)
            if ((co >= gridSize) or (ro >= gridSize)):
                print("Warning, shape outside of page area on page:{}".format(self.tag))
                return
            # Now - let's find out if there is a link...
            click_action = None
            if shape.shape_type != MSO_SHAPE_TYPE.GROUP:
                click_action = shape.click_action
            if click_action.hyperlink.address is not None:
                self.links[co][
                        ro] = click_action.hyperlink.address
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                if bordercolor:
                    # print "here"
                    # print "X
                    # {}".format(shape.line.color.rgb)
                    self.colors[co][
                            ro] = shape.line.color.rgb
                    # print "there"

                else:
                    if shape.fill.type==MSO_FILL.SOLID:
                        if (str(shape.fill.fore_color.type)
                                == "SCHEME (2)"):
                            #self.colors[co][ro] = ( 200, 0, 0)
                            pass
                        else:
                            self.colors[co][
                                    ro] = shape.fill.fore_color.rgb

            if shape.has_text_frame:
                self.process_text_frame(shape, co, ro)

            if(warningMissingLinks):
                if (click_action.hyperlink.address is None):
                    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                        if shape.auto_shape_type == MSO_SHAPE.FOLDED_CORNER:
                            if len(self.links[
                                   co][ro]) < 1:
                                self.pageset.addfeedback("Unknown link at slide: " + self.tag + " link here: [{}] [{}] {} ".format(co, ro, self.labels[co][ro]) + self.links[co][ro])

        except (AttributeError, KeyError, NotImplementedError):
            self.pageset.addfeedback(returnException())
            return

    def process_text_frame(self, shape, co, ro):
        text = ""
        if "Yes" in self.labels[co][ro]:
            return
        for paragraph in shape.text_frame.paragraphs:
            text += "".join([run.text.replace("’", "'")
                             for run in paragraph.runs])
        if text != "":
            # add the if shape_type is text box
            self.labels[co][ro] = text.strip()


def create_image_grid(slide, grid):
    images = {}
    for shape in slide.shapes:
        if not hasattr(shape, "shape_type"):
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            (co, ro) = grid.get_col_row(
                    shape.top+shape.height/2, shape.left+shape.width/2)
            if (co, ro) not in images:
                images[co, ro] = []
            images[co, ro].append(shape)
        if IMAGE_WARNING:
            for col in range(grid.grid_size):
                for row in range(grid.grid_size):
                    if (col, row) not in images:
                        if grid.labels[col][
                                row] is not "":
                            if grid.tag not in grid.labels[
                                    col][row]:
                                print("WARNING: image missing at column {}, row  {} (label: {}) on slide:{}".format(col, row, labels[col][row], grid.tag))
    return images


def export_images(grid, slide_number, slide, filename, SAVE=True):
    """     Second pass through shapes list finds images and saves them.
    We have to do this separately so it's guaranteed we already know what to
    name the images!"""
    labels = grid.labels
    images = create_image_grid(slide, grid)
    # Compose each icon out of all the images in the grid cell.
    for (x, y) in images:
        # Go through all the images, compute bounding
        # box.
        l = min([shape.left for shape in images[x, y]])
        t = min([shape.top for shape in images[x, y]])
        r = max([shape.left +
                 shape.width for shape in images[x, y]])
        b = max([shape.top +
                 shape.height for shape in images[x, y]])

        try:
            # Scale gives us the mapping from image pixels to powerpoint
            # distance units. This depends on the resolution
            # of the images.
            scale = min([shape.width/shape.image.size[0]
                         for shape in images[x, y]])

            # Size of combined image, in actual pixels (not PPTX units)
            # If scales differ between objects, we resize
            # them next
            w = int((r-l)/scale)
            h = int((b-t)/scale)

            composite = Image.new('RGBA', (w, h))

            # Add all the images together.
            for shape in images[x, y]:
                            # TODO: flipping.
                part = Image.open(
                    io.BytesIO(
                        shape.image.blob))
                part.load()
                width = part.size[0]
                height = part.size[1]
                left = shape.crop_left*width
                right = (1-shape.crop_right)*width
                top = shape.crop_top*height
                bottom = (1-shape.crop_bottom)*height
                box = (int(left),
                       int(top),
                       int(right),
                       int(bottom))
                part = part.crop(box)
                partScale = (shape.width / part.size[0])
                # part.size because it might have been cropped

                part = resizeImage(part, partScale / scale)
                composite.paste(
                    part,
                    (int((shape.left - l)/scale),int(
                     (shape.top - t)/scale)))
            # Crop final image.
            bbox = composite.getbbox()
            composite = composite.crop(bbox)

            # Save!
            grid.icons[x][y] = "icons/" + create_icon_name(x, y, labels, grid.links, slide_number)
            name =                  create_icon_name(x, y, labels, grid.links, slide_number)
            # print name
            if SAVE:
                # + str(slide_number)
                folder = filename+"/icons/"
                if not os.path.exists(folder):
                    os.makedirs(folder)
                composite.save(folder + "" + name)
        except IOError as e:
            print("Error reading image for {} {}".format(x, y))
            if ("cannot find loader for this WMF file" in e):
                print("Error: it appears that the image in column {} row {} of slide {}, is for Windows only, please change the format of that image".format(x, y, slide_number))
            if ("cannot identify image file" in e):
                print("Error: it appears that the image in column {} row {} of slide {}, is for Windows only, please change the format of that image".format(x, y, slide_number))
        except ValueError:
            print("Error reading image for {} {}".format(x, y))
        except IndexError:
            print("Error reading image for {} {}- it is outside the grid".format(x, y))



def create_json_object(grids):
    # Start the JSON output.
    grid_json = {}
    for i in range(len(grids)):
        grid_json[i] = [
            grids[i].tag,
            grids[i].labels,
            grids[i].utterances,
            grids[i].links,
            grids[i].icons,
            grids[i].colors,
            i]

    for_json = {}
    for_json["Settings"] = [gridSize, "test title", "en", ""]
    for_json["Grid"] = grid_json
    return for_json


def create_obf_manifest(root,boards_names_dic, image_names_dic, dest):
    # Create the manifest
#    root = boards_names_dic['toppage']
    string_of_board_names = json.dumps(boards_names_dic,ensure_ascii=False)
    string_of_image_names = json.dumps(image_names_dic,ensure_ascii=False)
    print("XXXXXXXXXXXXXXX")
    print(root)
    with io.open(dest+"/data/manifest.json", "w") as manifest:
        manifest.write("""{{
"format": "open-board-0.1",
"root": "{}",
"paths": {{
"boards":
{},
"images":
{}

}}
}}""".format(root, string_of_board_names, string_of_image_names))


def create_obf_button(grid,col,row):
    button = {}
    button["id"] = "{}{}".format(col, row)
    button["border_color"] = "rgb(68,68,68)"
    if("pptx" in str(type(grid.colors[col][row]))):
        color = grid.colors[col][row]
        button["background_color"] = "rgb({},{},{})".format(
            color[0], color[1], color[2])
    else:
        button["background_color"] = "rgb(0,0,0)"
    button["image_id"] = grid.icons[col][row]
    if len(grid.links[col][row]) > 1:
        if "special::" not in grid.links[col][row]:
            if not grid.links[col][row].startswith("ovf("):
                button["load_board"]= { "path": "boards/"+grid.links[col][row]+".obf" }
            else:
                #It's a special commend. let's extract it.
                print("We're in the special commands:")
                print("original line is : {}".format(grid.links[col][row]))
                commandstring=grid.links[col][row][4:-1]
                print("Command is {}".format(commandstring))
                commands=commandstring.split(",")
                for command in commands:
                    command_name= command.split("(",1)[0]
                    argument=command.split("(",1)[1][0:-1]
                    print("command name is {}".format(command_name))
                    print("argument is {}".format(argument))
                    if command_name == "deleteword":
                        #should find out what we do there...
                        button["action"]=":deleteword"
                    elif command_name == "clear":
                        button["action"]=":clear"
                    elif command_name == "place":
                        button["label"] = argument
                    elif command_name == "open":
                        button["load_board"]= { "path": "boards/"+make_title(argument)+".obf" }

                    elif command_name == "unfinnished":
                         pass
                    elif command_name == "blank":
                         pass
                    else:
                        raise ValueError('Unknown special command_name ({}) to process'.format(command_name))
                pass
    #This is at the end because we might change it during the special commands.
    button["label"] = grid.labels[col][row]
    return button



def create_obf_object(grid):
    for_json = {}
    for_json["format"] = "open-board-0.1"
    for_json["name"] = "CommuniKate "+make_title(grid.tag)
    for_json["locale"] = "en"
    for_json["id"] = make_title(grid.tag)
    for_json["grid"] = {}
    for_json["images"] = []
    for_json["sounds"] = []
    for_json["grid"]["rows"] = gridSize
    for_json["grid"]["columns"] = gridSize
    for_json["buttons"] = []
    for_json["grid"]["order"] = []
    for row in range(gridSize):
        grid_row = []
        for col in range(gridSize):
            if (len(grid.labels[col][row]) > 0):
                button=create_obf_button(grid,col,row)
                grid_row.append(button["id"])
                for_json["buttons"].append(button)
            else:
                grid_row.append(None)
        for_json["grid"]["order"].append(grid_row)
        for col in range(gridSize):
            if (len(grid.icons[col][row])>2):
                img = {}
                img["content_type"] = "image/png"
                img["id"] = grid.icons[col][row]
                img["width"] = 300
                img["height"] = 300
                img["path"]="../"+grid.icons[col][row]
                for_json["images"].append(img)
    return for_json


def write_to_obf(grids, dest):
    #Lots of relative addressing going on here.
    boards_names_dic = {}
    image_names_dic = {}
    owd = os.getcwd()
    root=make_title(grids[0].tag)
    print("The title of the root board is {}".format(root)) 
    for tok in grids:
        for_json = create_obf_object(tok)
        for image in for_json["images"]:
            image_names_dic[image['id']]=image['path']
        filename = 'boards/'+make_title(tok.tag)+'.obf'
        filename = filename #.encode('ascii', 'ignore') #so, this turns it into asci? 
        boards_names_dic[make_title(tok.tag)]=filename
        print(filename)
        with open(dest+'/data/'+filename, 'w') as outfile:
            json.dump(for_json, outfile, sort_keys=True, indent=2)
    create_obf_manifest(root,boards_names_dic,image_names_dic, dest)
    os.chdir(dest+'/data')
    outzipfile = 'pageset.obz'
    boards_names_dic['manifest']='manifest.json' #no idea what this line does, definately needs some test/reactoring.
    with zipfile.ZipFile(outzipfile, "w") as w:
        for x in list(boards_names_dic.values()):
            w.write(x)#.encode('utf8'))
        for y in list(image_names_dic.keys()):
            x=image_names_dic[y]
            image_names_dic[y]=x.replace("../icons","images")
            print(x)
            w.write(x) #.encode('utf8'))
    os.chdir(owd)




def write_to_JSON(grids, filename):
    for_json = create_json_object(grids)
    with open(filename, 'w') as outfile:
        json.dump(for_json, outfile, sort_keys=True, indent=4)



def make_title(label):
    """Given a  string, returns the string in the format we use for
     identifying grids. This is used mostly to build internal link
     structures"""
    tag = remove_punctuation(label.lower().strip().replace(" ", "_").replace("%20","_"))
    if tag == "":
        tag = "unknown"
    return tag

def create_ovf_manifest(filename):
    with open(filename, "w") as manifest:
        manifest.write("""{
"format": "open-board-0.1",
"root": "boards/my-board.ovf"
}""")

def create_icon_name(x, y, labels, links, slide_number):
    name = "S"+str(slide_number)+"X"+str(x)+"Y"+str(y)+".png" 
    try:
        name = "S"+str(slide_number)+"X"+str(x)+"Y"+str(y)+make_title(
                labels[x][y])+".png" 
    except IndexError:
        print("Create_icon_name was given an x y that was outside the possible ranges")
    return name

########

if __name__ == "__main__":
    if (len(sys.argv) < 2):
        print("\nUsage: ./grab_text.py <inputPptxFile> <gridSize>\n")
        print("inputPptxFile: The powerpoint pageset you want to process.")
        print("gridSize: width of square grid, e.g. '4' for a 4x4 grid")
        sys.exit(1)
    filename = sys.argv[1]
    dest = sys.argv[2]
    gridSize = 5
    if (len(sys.argv) > 2):
        gridSize = int(sys.argv[3])
    pageset = Pageset(filename, dest)
    write_to_JSON(pageset.grids, dest+'/pageset.json')
    write_to_JSON(pageset.grids, dest+'/pageset.json')
    create_ovf_manifest(dest+'/manifest.json')
    write_to_obf(pageset.grids, dest)

    # vim: tabstop=8 expandtab shiftwidth=4 softtabstop=4
