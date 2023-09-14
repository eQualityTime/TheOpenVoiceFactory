import math
import os
import io
import json
import urllib
import pagesetparser.core as core
import string
from PIL import Image
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL

warningMissingLinks = True
bordercolor = False #TODO - need a bordercolor regression test

class Grid:

    """Class representing on n by n grid, complete with utterances, links
    colours, and so on. """

    def __init__(self, pres, slide, grid_size, owningPageset,slide_number):
        self.button_order=[(i, j) for i in range(grid_size) for j in range(grid_size)]
        self.slide=slide
        self.slide_number=slide_number
        self.pageset = owningPageset
        self.grid_size = grid_size
        self.labels = [
            ["" for x in range(self.grid_size)]
            for x in range(self.grid_size)]
        self.links = [
            ["" for x in range(self.grid_size)]
            for x in range(self.grid_size)]
        self.colors = [
            ["" for x in range(self.grid_size)]
            for x in range(self.grid_size)]
        self.imagepaths = [
            ["" for x in range(self.grid_size)]
            for x in range(self.grid_size)]
        self.tag = "Slide number "+str(slide_number+1) #here for regression testing. 
        self.slide_width = pres.slide_width
        self.slide_height = pres.slide_height
        for shape in slide.shapes:
            self.process_shape(shape)

    @property
    def title(self):
        return core.make_title(self.tag) 

    def get_col_row_by_num(self, top: int, left: int):
        """
        Converts coordinates to column and row numbers based on grid size.
        Raises a ValueError if the shape is outside of the page area.
        """
        number_of_rows: int = self.grid_size
        number_of_columns: int = self.grid_size

        col: int = math.floor((number_of_columns * left / self.slide_width))
        row: int = math.floor((number_of_rows * top / self.slide_height))

        if col >= self.grid_size or row >= self.grid_size:
            raise ValueError("Shape outside of page area on page: {}".format(self.tag))

        return col, row


    def get_col_row(self,shape):
       return self.get_col_row_by_num(shape.top+shape.height/2, shape.left+shape.width/2)

    def process_shape(self, shape):
        try:
            if shape.is_placeholder:
                if shape.placeholder_format.idx == 0:
                    self.tag = core.make_title(shape.text)
            if shape.has_text_frame:  
                self.process_text_frame(shape)
            (co, ro) = self.get_col_row(shape)
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                if bordercolor:
                    self.colors[co][ro] = shape.line.color.rgb
                elif shape.fill.type==MSO_FILL.SOLID:
                    if (str(shape.fill.fore_color.type) != "SCHEME (2)"): 
                        self.colors[co][ro] = shape.fill.fore_color.rgb
            if shape.shape_type != MSO_SHAPE_TYPE.GROUP:
                click_action = shape.click_action
                if click_action.hyperlink.address is not None:
                    self.links[co][ro] = click_action.hyperlink.address

        except (ValueError, AttributeError, KeyError,  NotImplementedError):
            self.pageset.addfeedback(core.returnException())
            return

    def process_text_frame(self, shape):  
        text = ""
        (co,ro)=self.get_col_row(shape)
        if "Yes" in self.labels[co][ro]:  #This is a legacy hack - on one CK20 version (still used by some translations) there are multiple 'yes' labels on top of each other. Translators obviously only change the top one so this is here to make those translations work. 
            return
        for paragraph in shape.text_frame.paragraphs:
            text += "".join([run.text.replace("’", "'")
                             for run in paragraph.runs])  
        if text != "":
            self.labels[co][ro] = text.strip()

    def create_image_grid(self):
        images = {}
        for shape in self.slide.shapes:
            if not hasattr(shape, "shape_type"): #Regression tests don't trigger this
                continue 
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                (co, ro) = self.get_col_row(shape)
                images.setdefault((co,ro),[]).append(shape)
        return images

    def return_links(self):
        return_me=[]
        for (col,row) in self.button_order:
            return_me.append(self.links[row][col])
        return return_me

    def update_links(self, grids):  
        """the pptx file saves links to files 'slide1.xml', but we want them to point to the name of the boards."""
        for (col,row) in self.button_order:
            current = self.links[row][col]
            if "slide" in current:
                # first extract the number
                number= int(''.join( c for c in current if c in string.digits))
                # Then work out the relevant tag
                self.links[row][col] = grids[number - 1].tag
                # slides are numbered from 1 but grids are numbered from 0


    def write_obf_file(self,dest):
        filename = 'boards/'+self.title+'.obf'
        with open(dest+filename, 'w') as outfile:
            json.dump(self.create_obf_object(), outfile, sort_keys=True, indent=2)


    def create_obf_object(self):
        for_json = {}
        for_json["format"] = "open-board-0.1"
        for_json["name"] = "CommuniKate "+self.title #For V3 we should have the name be read from the pptx 
        for_json["locale"] = "en"
        for_json["id"] = self.title
        for_json["grid"] = {}
        for_json["images"] = []
        for_json["sounds"] = []
        for_json["grid"]["rows"] = self.grid_size
        for_json["grid"]["columns"] = self.grid_size
        for_json["buttons"] = []
        for_json["grid"]["order"] = []
        for row in range(self.grid_size):
            grid_row = []
            for col in range(self.grid_size):
                if (len(self.labels[col][row]) +len(self.links[col][row]) > 0): 
                    button=self.create_obf_button(col,row)
                    grid_row.append(button["id"])
                    for_json["buttons"].append(button)
                else:
                    grid_row.append(None)
            for_json["grid"]["order"].append(grid_row)

        for path in self.get_image_paths():
                    img = {}
                    img["content_type"] = "image/png" 
                    img["id"] = path
                    img["width"] = 300
                    img["height"] = 300
                    img["path"]=path
                    for_json["images"].append(img)
        return for_json


    def get_image_paths(self):
        image_paths=[]
        for row in range(self.grid_size):
            for col in range(self.grid_size):
#        for (col,row) in self.button_order: #TODO - put this back and regenerate the obz files (regression) - needs a full user test
                if (len(self.imagepaths[col][row])>2):
                    image_paths.append(self.imagepaths[col][row])
        return image_paths


    def create_obf_button(self,col,row): 
        button = {}
        button["id"] = "{}{}".format(col, row)
        button["label"] = self.labels[col][row]
        button["border_color"]= "rgb(68,68,68)"
        button["background_color"]="rgb(250,250,250)"
        if("pptx" in str(type(self.colors[col][row]))):
            button["background_color"] = "rgb({},{},{})".format( self.colors[col][row][0], self.colors[col][row][1], self.colors[col][row][2])
        button["image_id"] = self.imagepaths[col][row]
        if len(self.links[col][row]) > 1:
            if "special::" not in self.links[col][row]:  #TODO: Properly document this part @simple
                if not self.links[col][row].startswith("ovf("):
                    button["load_board"]= { "path": "boards/"+self.links[col][row]+".obf" }
                else:
                    core.process_commandstring(self.links[col][row],button)
        #This is at the end because we might change it during the special commands.
        if button["label"] is "": 
            button["label"]=button['id']
        return button

    def make_imagepaths(self):
        images = self.create_image_grid()
        for (x, y) in images: 
            self.imagepaths[x][y] = "images/" + self.icon_name(x, y)

    def icon_name(self, x, y):
        if 0 <= x < self.grid_size and 0 <= y < self.grid_size:
            pass #we've checked the inputs
        else:
            raise ValueError("Create_icon_name was given an x y that was outside the possible ranges")  
            return "Create_icon_name was given an x y that was outside the possible ranges"
            
        name = f"S{self.slide_number}X{x}Y{y}.png" 
        try:
            name = "S"+str(self.slide_number)+"X"+str(x)+"Y"+str(y)+core.make_title(self.labels[x][y])+".png" 
        except IndexError: #The IF statement at the top should mean this is never triggered (it isn't in the regression tests) 
            print("ERROR - Create_icon_name has raised an indexError") 
        return name


    def export_images(self, dest_folder, SAVE=True): 
        if SAVE:  
            folder = dest_folder+"/images/"
            if not os.path.exists(folder):
                os.makedirs(folder)
        images = self.create_image_grid()
        for (x, y) in images: 
            # Go through all the images, compute bounding box.
            left = min([shape.left for shape in images[x, y]]) 
            top = min([shape.top for shape in images[x, y]])
            right = max([shape.left + shape.width for shape in images[x, y]])
            bottom = max([shape.top + shape.height for shape in images[x, y]])
            # Scale gives us the mapping from image pixels to powerpoint distance units. This depends on the resolution of the images.
            scale = min([shape.width/shape.image.size[0] for shape in images[x, y]])
            # Size of combined image, in actual pixels (not PPTX units)
            # If scales differ between objects, we resize them next
            width = int((right-left)/scale)
            height = int((bottom-top)/scale)
            try:
                composite = Image.new('RGBA', (width, height))
                # Add all the images together.
                for shape in images[x, y]:
                    partBox=ready_for_composite(shape,scale,width,height,left,top)
                    composite=Image.alpha_composite(composite,partBox)
                # Crop final image.
                bbox = composite.getbbox()
                composite = composite.crop(bbox)
                # Save!
                self.imagepaths[x][y] = "images/" +self.icon_name(x, y)
                if SAVE:  
                    composite.save(dest_folder +"/"+ self.imagepaths[x][y])
            except IOError as e:
                print("Error reading image for {} {} (Code 121)".format(x, y))
                if ("cannot find loader for this WMF file" in str(e)):
                    print("Error: it appears that the image in column {} row {} of slide {}, is for Windows only, please change the format of that image".format(x, y, self.slide_number))
                if ("cannot identify image file" in str(e)):
                    print("Error: it appears that the image in column {} row {} of slide {}, is for Windows only, please change the format of that image".format(x, y, self.slide_number))
            except ValueError:
                print("Error reading image for {} {} (Code 127)".format(x, y))
            except IndexError:
                print("Error reading image for {} {}- it is outside the grid".format(x, y))


 #TODO: work out where these functions should be
# Options: 
# * A image module 
# * An image processing class 
# * The grid class 
def ready_for_composite(shape,scale,w,h,l,t):

    part = Image.open(
        io.BytesIO(
            shape.image.blob))
    part.load()
    part = part.crop(crop_to_shape(shape,part))
    partScale = (shape.width / part.size[0])
    # part.size because it might have been cropped
    part = resizeImage(part, partScale / scale)
    partBox = Image.new('RGBA', (w, h))
    partBox.paste( part, (int((shape.left - l)/scale),int( (shape.top - t)/scale)))
    return partBox

def crop_to_shape(shape,part):
    width = part.size[0]
    height = part.size[1]
    left = shape.crop_left*width
    right = (1-shape.crop_right)*width
    top = shape.crop_top*height
    bottom = (1-shape.crop_bottom)*height
    box = (int(left),
           int(top),
           int(right),
           int(bottom))
    return box


def resizeImage(image, scaleFactor):
    oldSize = image.size
    newSize = (int(scaleFactor*oldSize[0]),
               int(scaleFactor*oldSize[1]))
    return image.resize(newSize, Image.LANCZOS)




